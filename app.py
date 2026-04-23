import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import streamlit as st
import sys

st.write(sys.executable)

# ── Color palette ──────────────────────────────────────────────────────────
DEEP_NAVY  = RGBColor(0x0D, 0x1B, 0x2A)
NAVY       = RGBColor(0x1A, 0x2B, 0x3C)
TEAL       = RGBColor(0x00, 0xC6, 0xFF)
BLUE       = RGBColor(0x00, 0x72, 0xFF)
CYAN       = RGBColor(0x00, 0xE5, 0xFF)
DARK_CARD  = RGBColor(0x11, 0x22, 0x33)
CARD_BORDER= RGBColor(0x1E, 0x3A, 0x5F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xC4, 0xD8)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
RED        = RGBColor(0xFF, 0x52, 0x52)
YELLOW     = RGBColor(0xFF, 0xD7, 0x40)
ACCENT     = RGBColor(0x00, 0xBC, 0xD4)

W = Inches(10)      # slide width
H = Inches(5.625)   # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if line_color is None else None
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = font_name
    return txb

def mpl_to_pptx_image(fig):
    """Convert a matplotlib figure to an in-memory PNG stream."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor=fig.get_facecolor(), dpi=150)
    buf.seek(0)
    plt.close(fig)
    return buf

def add_oval(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Login & Register
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, DEEP_NAVY)

# Left panel
add_rect(s1, 0, 0, 4.5, 5.625, fill_color=NAVY)

# Decorative signal circles
for i in range(4):
    r = RGBColor(0x00, 0xC6, 0xFF)
    add_oval(s1, 0.3 - i*0.12, 0.2 - i*0.12,
             1.2 + i*0.5, 1.2 + i*0.5, fill_color=NAVY, line_color=TEAL, line_width=Pt(0.5))

# Brand
add_text(s1, "TELECOM", 0.15, 1.05, 4.2, 0.55,
         font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(s1, "CHURN PREDICTOR", 0.15, 1.55, 4.2, 0.35,
         font_size=11, color=TEAL, align=PP_ALIGN.CENTER)

# Horizontal lines
for i in range(5):
    add_rect(s1, 0.5, 2.0 + i*0.28, 3.5, 0.02, fill_color=TEAL)

# Stats cards
stats = [("98.5%", "Network Uptime"), ("4.2M", "Subscribers"), ("AI", "Powered")]
for i, (val, label) in enumerate(stats):
    y = 3.3 + i * 0.68
    add_rect(s1, 0.3, y, 3.9, 0.55, fill_color=DARK_CARD, line_color=TEAL, line_width=Pt(0.8))
    add_text(s1, val,   0.4,  y+0.05, 1.0, 0.45, font_size=14, bold=True, color=TEAL, font_name="Arial Black")
    add_text(s1, label, 1.45, y+0.08, 2.7, 0.40, font_size=11, color=LIGHT_GRAY)

# Right: Login
add_text(s1, "Welcome Back", 4.8, 0.5, 4.8, 0.55,
         font_size=26, bold=True, color=WHITE, font_name="Arial Black")
add_text(s1, "Sign in to access your dashboard", 4.8, 1.05, 4.8, 0.3,
         font_size=11, color=LIGHT_GRAY)

for label, y in [("Email Address", 1.55), ("Password", 2.35)]:
    add_text(s1, label, 4.8, y, 4.8, 0.28, font_size=10, color=LIGHT_GRAY)
    add_rect(s1, 4.8, y+0.3, 4.8, 0.42, fill_color=DARK_CARD, line_color=TEAL, line_width=Pt(1))

add_rect(s1, 4.8, 3.0, 4.8, 0.5, fill_color=BLUE)
add_text(s1, "SIGN IN", 4.8, 3.0, 4.8, 0.5,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s1, "— OR —", 4.8, 3.6, 4.8, 0.3, font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s1, "New User? Register Here", 4.8, 3.95, 4.8, 0.28,
         font_size=11, bold=True, color=WHITE)
add_rect(s1, 4.8, 4.28, 4.8, 0.5, fill_color=DEEP_NAVY, line_color=GREEN, line_width=Pt(1.5))
add_text(s1, "CREATE ACCOUNT", 4.8, 4.28, 4.8, 0.5,
         font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s1, "Slide 1 of 5  |  Telecom Churn Prediction System",
         0, 5.35, 10, 0.27, font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Customer Information Input
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, DEEP_NAVY)

add_rect(s2, 0, 0, 10, 0.8, fill_color=NAVY)
add_rect(s2, 0, 0, 0.05, 5.625, fill_color=TEAL)
add_text(s2, "Customer Information Input", 0.75, 0.18, 7.5, 0.45,
         font_size=20, bold=True, color=WHITE, font_name="Arial Black")
add_text(s2, "Step 2 of 5", 8.8, 0.22, 1.1, 0.36,
         font_size=10, color=TEAL, align=PP_ALIGN.CENTER)

# Progress bar
add_rect(s2, 0.15, 0.82, 9.7, 0.1, fill_color=CARD_BORDER)
add_rect(s2, 0.15, 0.82, 3.88, 0.1, fill_color=TEAL)

left_fields = [
    ("Call Failure",        "0 – 100",     RED),
    ("Subscription Length", "1 – 60 months", TEAL),
    ("Seconds of Use",      "0 – 100,000", BLUE),
    ("Age",                 "18 – 80",     ACCENT),
    ("Customer Value",      "0 – 1,000",   YELLOW),
]
right_fields = [
    ("Complains",           "Yes / No",    RED),
    ("Charge Amount",       "0 – 1,000",   GREEN),
    ("Frequency of Use",    "0 – 1,000",   TEAL),
    ("Frequency of SMS",    "0 – 1,000",   CYAN),
    ("Cluster",             "Low / Mid / High", YELLOW),
]

for i, (label, hint, accent) in enumerate(left_fields):
    y = 1.1 + i * 0.82
    add_rect(s2, 0.15, y, 4.55, 0.65, fill_color=DARK_CARD, line_color=accent, line_width=Pt(0.8))
    add_rect(s2, 0.15, y, 0.05, 0.65, fill_color=accent)
    add_text(s2, label, 0.3, y+0.05, 4.3, 0.25, font_size=10, bold=True, color=WHITE)
    add_text(s2, hint,  0.3, y+0.32, 4.3, 0.25, font_size=9,  color=LIGHT_GRAY)

for i, (label, hint, accent) in enumerate(right_fields):
    y = 1.1 + i * 0.82
    add_rect(s2, 5.0, y, 4.85, 0.65, fill_color=DARK_CARD, line_color=accent, line_width=Pt(0.8))
    add_rect(s2, 5.0, y, 0.05, 0.65, fill_color=accent)
    add_text(s2, label, 5.15, y+0.05, 4.55, 0.25, font_size=10, bold=True, color=WHITE)
    add_text(s2, hint,  5.15, y+0.32, 4.55, 0.25, font_size=9,  color=LIGHT_GRAY)

bottom = [("Tariff Plan","Basic / Premium",TEAL),
          ("Status","Active / Inactive",GREEN),
          ("Distinct Called Numbers","0 – 200",YELLOW)]
for i, (label, hint, accent) in enumerate(bottom):
    x = 0.15 + i * 3.28
    add_rect(s2, x, 5.13, 3.1, 0.35, fill_color=DARK_CARD, line_color=accent, line_width=Pt(0.8))
    add_text(s2, f"{label}  ·  {hint}", x+0.1, 5.16, 2.9, 0.3, font_size=8.5, color=LIGHT_GRAY)

add_text(s2, "Slide 2 of 5  |  Customer Data Entry",
         0, 5.5, 10, 0.12, font_size=7.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Prediction + AI Explanation
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3, DEEP_NAVY)

add_rect(s3, 0, 0, 10, 0.75, fill_color=NAVY)
add_rect(s3, 0, 0, 10, 0.04, fill_color=TEAL)
add_text(s3, "Prediction Result & AI Explanation", 0.75, 0.18, 7.5, 0.42,
         font_size=18, bold=True, color=WHITE, font_name="Arial Black")
add_text(s3, "Step 3 of 5", 8.8, 0.2, 1.1, 0.36, font_size=10, color=TEAL, align=PP_ALIGN.CENTER)

# Churn result card
add_rect(s3, 0.15, 0.9, 4.5, 1.6,
         fill_color=RGBColor(0x1A, 0x0A, 0x0A), line_color=RED, line_width=Pt(2))
add_text(s3, "CHURN PREDICTED", 1.05, 1.0, 3.5, 0.45,
         font_size=15, bold=True, color=RED, font_name="Arial Black")
add_text(s3, "Customer is likely to leave the network",
         0.35, 1.5, 4.2, 0.3, font_size=10, color=LIGHT_GRAY)
add_rect(s3, 0.35, 1.85, 4.1, 0.12, fill_color=CARD_BORDER)
add_rect(s3, 0.35, 1.85, 3.2,  0.12, fill_color=RED)
add_text(s3, "Churn Probability: 78%", 0.35, 2.02, 4.1, 0.3,
         font_size=10, bold=True, color=RED)

# Risk card
add_rect(s3, 4.85, 0.9, 4.9, 1.6, fill_color=DARK_CARD, line_color=YELLOW, line_width=Pt(1.5))
add_text(s3, "Risk Level: HIGH", 5.0, 1.0, 4.55, 0.38,
         font_size=14, bold=True, color=YELLOW, font_name="Arial Black")
risk_items = ["High call failure rate (>5)",
              "Short subscription (<12 months)",
              "Low usage frequency (<10)"]
for i, txt in enumerate(risk_items):
    add_oval(s3, 5.0, 1.44+i*0.32, 0.12, 0.12, fill_color=RED)
    add_text(s3, txt, 5.18, 1.42+i*0.32, 4.35, 0.28, font_size=9.5, color=LIGHT_GRAY)

# ── Bar chart: Feature Impact ─────────────────────────────────────────────
add_text(s3, "Feature Impact on Churn", 0.15, 2.65, 4.5, 0.3,
         font_size=11, bold=True, color=TEAL)

fig_bar, ax = plt.subplots(figsize=(4.5, 2.5))
fig_bar.patch.set_facecolor("#112233")
ax.set_facecolor("#112233")
labels = ["Call\nFailure", "Complains", "Sub\nLength", "Cust\nValue", "Usage\nFreq"]
values = [85, 78, 62, 55, 40]
bar_colors = ["#FF5252","#FF5252","#FFD740","#FFD740","#00C6FF"]
bars = ax.bar(labels, values, color=bar_colors, width=0.6)
for bar, v in zip(bars, values):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
            str(v), ha="center", va="bottom", color="white", fontsize=9, fontweight="bold")
ax.set_ylim(0, 100)
ax.tick_params(colors="#B0C4D8", labelsize=8)
ax.spines[:].set_visible(False)
ax.yaxis.set_visible(False)
for spine in ax.spines.values():
    spine.set_edgecolor("#1E3A5F")
plt.tight_layout(pad=0.3)
s3.shapes.add_picture(mpl_to_pptx_image(fig_bar), Inches(0.15), Inches(2.95), Inches(4.5), Inches(2.3))

# ── Doughnut chart: Segment Risk ─────────────────────────────────────────
add_text(s3, "Customer Segment Risk", 5.0, 2.65, 4.9, 0.3,
         font_size=11, bold=True, color=TEAL)

fig_d, ax2 = plt.subplots(figsize=(4.5, 2.5))
fig_d.patch.set_facecolor("#112233")
ax2.set_facecolor("#112233")
sizes = [32, 45, 23]
seg_labels = ["High Risk", "Medium Risk", "Low Risk"]
seg_colors = ["#FF5252", "#FFD740", "#00E676"]
wedges, texts, autotexts = ax2.pie(
    sizes, labels=seg_labels, colors=seg_colors,
    autopct="%1.0f%%", startangle=90,
    pctdistance=0.75, wedgeprops=dict(width=0.5))
for t in texts:      t.set_color("#B0C4D8"); t.set_fontsize(8)
for t in autotexts:  t.set_color("white");   t.set_fontsize(8); t.set_fontweight("bold")
plt.tight_layout(pad=0.3)
s3.shapes.add_picture(mpl_to_pptx_image(fig_d), Inches(5.0), Inches(2.95), Inches(4.9), Inches(2.3))

add_text(s3, "Slide 3 of 5  |  AI-Powered Prediction Engine",
         0, 5.5, 10, 0.12, font_size=7.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Download CSV
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4, DEEP_NAVY)

add_rect(s4, 0, 0, 10, 0.75, fill_color=NAVY)
add_rect(s4, 0, 0, 10, 0.04, fill_color=GREEN)
add_text(s4, "Download Results & Analytics", 0.75, 0.18, 7.5, 0.42,
         font_size=18, bold=True, color=WHITE, font_name="Arial Black")
add_text(s4, "Step 4 of 5", 8.8, 0.2, 1.1, 0.36, font_size=10, color=GREEN, align=PP_ALIGN.CENTER)

# CSV preview card
add_rect(s4, 0.15, 0.88, 5.7, 3.0, fill_color=DARK_CARD, line_color=GREEN, line_width=Pt(1.5))
add_text(s4, "churn_result.csv", 0.35, 0.98, 5.3, 0.35,
         font_size=12, bold=True, color=GREEN)

# Mini table
rows = [("Call Failure","12",LIGHT_GRAY,CYAN),
        ("Complains","Yes",LIGHT_GRAY,CYAN),
        ("Sub. Length","8 months",LIGHT_GRAY,CYAN),
        ("Customer Value","145",LIGHT_GRAY,CYAN),
        ("Prediction","CHURN",LIGHT_GRAY,RED),
        ("Probability","0.78",LIGHT_GRAY,CYAN)]
add_rect(s4, 0.35, 1.38, 5.3, 0.3, fill_color=RGBColor(0x0D,0x20,0x40))
add_text(s4, "Field", 0.4, 1.4, 2.8, 0.26, font_size=9, bold=True, color=WHITE)
add_text(s4, "Value", 3.2, 1.4, 2.2, 0.26, font_size=9, bold=True, color=WHITE)
for i, (f, v, fc, vc) in enumerate(rows):
    y = 1.72 + i * 0.29
    add_rect(s4, 0.35, y, 5.3, 0.29,
             fill_color=DARK_CARD if i%2==0 else RGBColor(0x0D,0x1B,0x2A))
    add_text(s4, f, 0.4,  y+0.03, 2.7, 0.24, font_size=9, color=fc)
    add_text(s4, v, 3.2,  y+0.03, 2.2, 0.24, font_size=9, color=vc)

# Download button
add_rect(s4, 0.15, 4.05, 5.7, 0.55,
         fill_color=RGBColor(0x0A,0x3A,0x1A), line_color=GREEN, line_width=Pt(2))
add_text(s4, "DOWNLOAD CSV REPORT", 0.15, 4.05, 5.7, 0.55,
         font_size=12, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, font_name="Arial Black")

# ── Line chart: churn trend ───────────────────────────────────────────────
add_text(s4, "Session Analytics", 6.1, 0.88, 3.75, 0.32, font_size=11, bold=True, color=TEAL)
fig_line, ax3 = plt.subplots(figsize=(3.75, 1.8))
fig_line.patch.set_facecolor("#112233")
ax3.set_facecolor("#112233")
months = ["Jan","Feb","Mar","Apr","May","Jun"]
probs  = [0.42, 0.55, 0.61, 0.70, 0.75, 0.78]
ax3.plot(months, probs, color="#FF5252", linewidth=2.5, marker="o", markersize=5)
ax3.fill_between(months, probs, alpha=0.15, color="#FF5252")
ax3.tick_params(colors="#B0C4D8", labelsize=7)
ax3.spines["top"].set_visible(False); ax3.spines["right"].set_visible(False)
for sp in ["bottom","left"]: ax3.spines[sp].set_edgecolor("#1E3A5F")
ax3.set_ylim(0, 1)
plt.tight_layout(pad=0.3)
s4.shapes.add_picture(mpl_to_pptx_image(fig_line), Inches(6.1), Inches(1.2), Inches(3.75), Inches(1.8))

# ── Pie chart: prediction distribution ───────────────────────────────────
add_text(s4, "Prediction Distribution", 6.1, 3.1, 3.75, 0.3, font_size=11, bold=True, color=TEAL)
fig_pie, ax4 = plt.subplots(figsize=(3.75, 1.8))
fig_pie.patch.set_facecolor("#112233")
ax4.set_facecolor("#112233")
ax4.pie([78, 22], labels=["Will Churn","Won't Churn"],
        colors=["#FF5252","#00E676"], autopct="%1.0f%%",
        textprops={"color":"#B0C4D8","fontsize":8},
        pctdistance=0.75, startangle=90)
plt.tight_layout(pad=0.3)
s4.shapes.add_picture(mpl_to_pptx_image(fig_pie), Inches(6.1), Inches(3.42), Inches(3.75), Inches(1.8))

add_text(s4, "Slide 4 of 5  |  Export & Download Center",
         0, 5.5, 10, 0.12, font_size=7.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Logout
# ══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5, DEEP_NAVY)

# Decorative radar circles
for i in range(5):
    add_oval(s5, 5.0-i*0.6, 2.5-i*0.6, 1.5+i*1.2, 1.5+i*1.2,
             fill_color=DEEP_NAVY, line_color=TEAL, line_width=Pt(0.5+i*0.2))

# Center panel
add_rect(s5, 2.5, 0.8, 5.0, 4.0, fill_color=NAVY, line_color=TEAL, line_width=Pt(1))

# Check circle
add_oval(s5, 4.5, 1.0, 1.0, 1.0,
         fill_color=RGBColor(0x00,0x22,0x11), line_color=GREEN, line_width=Pt(2))
add_text(s5, "✓", 4.5, 1.0, 1.0, 1.0,
         font_size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s5, "Session Complete!", 2.5, 2.1, 5.0, 0.55,
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(s5, "Your prediction results have been saved.",
         2.6, 2.72, 4.8, 0.3, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Summary stats
summary = [("1","Predictions Run",TEAL), ("78%","Churn Probability",RED), ("1","CSV Exported",GREEN)]
for i, (val, label, col) in enumerate(summary):
    x = 2.65 + i * 1.6
    add_rect(s5, x, 3.1, 1.45, 0.8, fill_color=DARK_CARD, line_color=col, line_width=Pt(1))
    add_text(s5, val,   x, 3.12, 1.45, 0.38, font_size=16, bold=True,
             color=col, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(s5, label, x, 3.50, 1.45, 0.35, font_size=8,
             color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Logout button
add_rect(s5, 3.25, 4.05, 3.5, 0.52,
         fill_color=RGBColor(0x2A,0x00,0x00), line_color=RED, line_width=Pt(2))
add_text(s5, "LOGOUT", 3.25, 4.05, 3.5, 0.52,
         font_size=13, bold=True, color=RED,
         align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s5, "TELECOM CHURN PREDICTOR  ·  Powered by AI",
         0, 5.28, 10, 0.3, font_size=9, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s5, "Slide 5 of 5  |  Session Terminated",
         0, 5.5, 10, 0.12, font_size=7.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════
prs.save("telecom_churn_presentation.pptx")
print("Saved: telecom_churn_presentation.pptx")
st.markdown("---")
st.markdown("🚀 Telecom Churn Prediction | BCA Project")
